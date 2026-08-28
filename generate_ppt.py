from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(11, 2, 22)

PINK = RGBColor(249, 168, 212)
PURPLE = RGBColor(167, 139, 250)
WHITE = RGBColor(255, 255, 255)
DIM = RGBColor(161, 161, 170)
HOTPINK = RGBColor(244, 114, 182)

def add_text(tf, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_after=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_after = space_after
    return p

# ─── TITLE ───
tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Myntra Mindsight"
p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
add_text(tf, "AI-Powered Fashion Purchase Discovery Intelligence Engine", 20, DIM)

# Date tag
tb2 = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.6))
tf2 = tb2.text_frame
p = tf2.paragraphs[0]
p.text = "INTERNAL PRESENTATION"; p.font.size = Pt(9); p.font.color.rgb = PURPLE; p.font.bold = True; p.alignment = PP_ALIGN.RIGHT
add_text(tf2, "August 2026", 9, DIM, align=PP_ALIGN.RIGHT)

# ─── PIPELINE STRIP ───
tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(0.8))
tf3 = tb3.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "HOW IT WORKS — AI ANALYSIS PIPELINE"; p.font.size = Pt(9); p.font.color.rgb = PURPLE; p.font.bold = True
add_text(tf3, "Ingest 804 Reviews  →  Clean & De-spam (156 removed)  →  Groq LLM Enrichment  →  17 AI Signals Extracted  →  Semantic Clustering (ChromaDB)  →  8 Ranked Opportunities  →  PM Dashboard & Chat", 13, RGBColor(228, 228, 231))

# ─── LEFT COLUMN ───
left_x = Inches(0.6)
col_w = Inches(5.8)

# Stats row
stats = [("804", "Feedback Items"), ("7", "Data Sources"), ("17", "AI Signals"), ("8", "Opportunities"), ("648", "Clean Reviews")]
stat_colors = [WHITE, RGBColor(192, 132, 252), HOTPINK, RGBColor(56, 189, 248), RGBColor(74, 222, 128)]
for i, (num, label) in enumerate(stats):
    x = left_x + Inches(i * 1.16)
    tb = slide.shapes.add_textbox(x, Inches(2.55), Inches(1.05), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = stat_colors[i]; p.alignment = PP_ALIGN.CENTER
    add_text(tf, label, 8, DIM, align=PP_ALIGN.CENTER)

# Source Platform Breakdown
tb = slide.shapes.add_textbox(left_x, Inches(3.5), col_w, Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "SOURCE PLATFORM BREAKDOWN"; p.font.size = Pt(9); p.font.color.rgb = PINK; p.font.bold = True
add_text(tf, "Play Store: 386  |  Reddit: 334  |  Product Reviews: 50  |  App Store: 10  |  Instagram: 8  |  Twitter/X: 8  |  YouTube: 8", 13, RGBColor(212,212,216))

# Journey Stage Distribution
tb = slide.shapes.add_textbox(left_x, Inches(4.4), col_w, Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "SIGNAL DISTRIBUTION BY JOURNEY STAGE"; p.font.size = Pt(9); p.font.color.rgb = PINK; p.font.bold = True
add_text(tf, "Evaluation: 5  |  Post-Purchase: 3  |  Product Consideration: 3  |  Wishlist: 3  |  Comparison: 2  |  Purchase Postponement: 1", 13, RGBColor(212,212,216))

# AI Signal Types
tb = slide.shapes.add_textbox(left_x, Inches(5.3), col_w, Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "AI SIGNAL TYPES DETECTED"; p.font.size = Pt(9); p.font.color.rgb = PINK; p.font.bold = True
add_text(tf, "Barriers: 7  |  Unmet Needs: 7  |  Social Validation: 1  |  Risk Avoidance: 1  |  Avoidance Behavior: 1", 13, RGBColor(212,212,216))

# ─── RIGHT COLUMN: Opportunities ───
right_x = Inches(6.8)
right_w = Inches(6.2)

tb = slide.shapes.add_textbox(right_x, Inches(2.4), right_w, Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "TOP 8 RANKED OPPORTUNITIES (BY COMPOSITE SCORE)"; p.font.size = Pt(9); p.font.color.rgb = PINK; p.font.bold = True

import sqlite3
try:
    conn = sqlite3.connect('backend/data/discovery_engine.db')
    cursor = conn.cursor()
    cursor.execute('''
        SELECT o.title, o.description, o.supporting_conversations, os.composite_score
        FROM opportunity o
        JOIN opportunity_score os ON o.id = os.opportunity_id
        ORDER BY os.composite_score DESC
        LIMIT 8
    ''')
    rows = cursor.fetchall()
    conn.close()
    
    colors = [
        HOTPINK,
        RGBColor(192,132,252),
        RGBColor(56,189,248),
        DIM,
        DIM,
        DIM,
        DIM,
        DIM
    ]
    
    opps = []
    for i, row in enumerate(rows):
        rank = f"#{i+1}"
        title = row[0]
        # Shorten description to fit PPT
        desc_text = (row[1][:90] + '...') if len(row[1]) > 90 else row[1]
        desc_full = f"{desc_text} · {row[2]} conv."
        score = f"{row[3]:.1f}/100"
        opps.append((rank, title, desc_full, score, colors[i]))
        
except Exception as e:
    print("Could not load from DB:", e)
    opps = []

for i, (rank, title, desc, score, color) in enumerate(opps):
    y = Inches(2.8) + Inches(i * 0.46)
    tb = slide.shapes.add_textbox(right_x, y, right_w, Inches(0.44))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{rank}  {title}   —   Score: {score}"
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = color
    add_text(tf, desc, 10, DIM)

# ─── CORE FINDING ───
tb = slide.shapes.add_textbox(right_x, Inches(6.55), right_w, Inches(0.7))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "CORE FINDING"; p.font.size = Pt(10); p.font.color.rgb = PINK; p.font.bold = True
add_text(tf, "7 of 17 AI signals are purchase barriers and 7 are unmet needs — concentrated in Evaluation & Post-Purchase stages. All 17 signals have ≥ 0.90 confidence. The #1 opportunity has the highest reach (0.30) and widest evidence (6 conversations).", 11, RGBColor(212,212,216))

# ─── FOOTER ───
tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(6), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Tech: FastAPI · SQLAlchemy · SQLite · ChromaDB · Groq LLM · Next.js · React"; p.font.size = Pt(9); p.font.color.rgb = RGBColor(82,82,91)

tb = slide.shapes.add_textbox(Inches(7), Inches(7.1), Inches(6), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Pipeline: Question → Evidence → Analysis → Pattern → Opportunity → Business Relevance"; p.font.size = Pt(9); p.font.color.rgb = RGBColor(82,82,91); p.alignment = PP_ALIGN.RIGHT

prs.save("Myntra_Mindsight.pptx")
print("✅ Myntra_Mindsight.pptx created successfully")
