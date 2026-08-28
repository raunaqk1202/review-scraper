import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 6 is blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(21, 14, 33) # Deep dark purple to match reference exactly
    
    def add_text_box(left, top, width, height, text, font_size=14, bold=False, color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.name = 'Calibri'
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        return txBox

    # TITLE
    add_text_box(Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "Myntra Mindsight", font_size=44, bold=True, color=RGBColor(255, 255, 255))
    
    # SUBTITLE
    add_text_box(Inches(0.5), Inches(1.1), Inches(12), Inches(0.4), "AI-Powered Fashion Purchase Discovery Intelligence Engine", font_size=20, bold=False, color=RGBColor(160, 160, 170))

    heading_color = RGBColor(180, 130, 240)

    # HOW IT WORKS
    add_text_box(Inches(0.5), Inches(1.7), Inches(6), Inches(0.3), "HOW IT WORKS — AI ANALYSIS PIPELINE", font_size=14, bold=True, color=heading_color)
    
    pipeline_text = "Ingest 804 Reviews  →  Clean & De-spam (156 removed)  →  Groq LLM Enrichment  →  17 AI Signals Extracted  →  Semantic Clustering (ChromaDB)  →  8 Ranked Opportunities  →  PM Dashboard & Chat"
    add_text_box(Inches(0.5), Inches(2.0), Inches(6.2), Inches(0.6), pipeline_text, font_size=14, color=RGBColor(230, 230, 230))
    
    # NUMBERS SECTION
    numbers_top = 2.7
    stats = [
        ("804", "Feedback\nItems", RGBColor(255, 255, 255)),
        ("7", "Data\nSources", RGBColor(180, 130, 240)),
        ("648", "Clean\nReviews", RGBColor(100, 220, 120)),
        ("17", "AI Signals", RGBColor(255, 105, 180)),
        ("8", "Opportunities", RGBColor(100, 180, 255))
    ]
    
    left_offset = 0.5
    for val, label, color in stats:
        # Value
        tx = add_text_box(Inches(left_offset), Inches(numbers_top), Inches(1.0), Inches(0.5), val, font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Label
        lbl = add_text_box(Inches(left_offset), Inches(numbers_top + 0.5), Inches(1.0), Inches(0.5), label, font_size=12, color=RGBColor(180, 180, 180), align=PP_ALIGN.CENTER)
        left_offset += 1.2
        
    # DETAILED HORIZONTAL FLOWCHART
    add_text_box(Inches(0.5), Inches(3.6), Inches(6), Inches(0.3), "HOW IT WORKS (DETAILED PIPELINE)", font_size=12, bold=True, color=heading_color)
    
    steps = [
        ("1. Data Ingestion", "Scrapes Play Store, App Store & Reddit. Cleans spam & normalizes raw user reviews."),
        ("2. AI Processing", "Groq LLM extracts Pain Points, Needs, & Barriers from unstructured text."),
        ("3. Clustering", "Semantic search groups similar signals into recurring themes & patterns."),
        ("4. Opp Generation", "Generates Product Opportunities & ranks them using the 4-dimension scoring model.")
    ]
    
    box_w = 1.45
    box_h = 2.0
    start_x = 0.5
    start_y = 4.0
    
    for i, (title, desc) in enumerate(steps):
        # Draw Box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(start_y), Inches(box_w), Inches(box_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(35, 20, 50)
        shape.line.color.rgb = RGBColor(180, 130, 240)
        shape.line.width = Pt(1)
        
        # Add Title
        tf = shape.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.1)
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add Description
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = 'Calibri'
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(200, 200, 200)
        
        # Draw Arrow
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x + box_w + 0.02), Inches(start_y + (box_h / 2) - 0.1), Inches(0.16), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(180, 130, 240)
            arrow.line.color.rgb = RGBColor(180, 130, 240)
            
        start_x += box_w + 0.2

    # SCORING ALGORITHM & ASSUMPTIONS 
    add_text_box(Inches(0.5), Inches(6.2), Inches(6), Inches(0.3), "SCORING ALGORITHM & ASSUMPTIONS", font_size=10, bold=True, color=heading_color)
    add_text_box(Inches(0.5), Inches(6.4), Inches(6.2), Inches(0.8), "Formula: 35% Pain + 30% Impact + 20% Reach + 15% Evidence\nAssumptions: LLM scales 1.0-5.0. User Pain is primary driver. Avoids indexing single anecdotal claims. Scaled to 0-100 composite.", font_size=12, color=RGBColor(230, 230, 230))

    # TOP 8 RANKED OPPORTUNITIES (Right side)
    add_text_box(Inches(6.8), Inches(2.4), Inches(6), Inches(0.3), "TOP 8 RANKED OPPORTUNITIES (BY COMPOSITE SCORE)", font_size=14, bold=True, color=heading_color)
    
    import sqlite3
    try:
        conn = sqlite3.connect('backend/data/discovery_engine.db')
        cursor = conn.cursor()
        cursor.execute('''
            SELECT o.title, o.description, o.supporting_conversations, os.composite_score
            FROM opportunity o
            JOIN opportunity_score os ON o.id = os.opportunity_id
            ORDER BY os.composite_score DESC
            LIMIT 8
        ''')
        rows = cursor.fetchall()
        conn.close()
        
        colors = [
            RGBColor(255, 105, 180),
            RGBColor(180, 130, 240),
            RGBColor(100, 180, 255),
            RGBColor(100, 220, 120),
            RGBColor(255, 255, 255),
            RGBColor(255, 200, 50),
            RGBColor(255, 150, 50),
            RGBColor(255, 80, 80)
        ]
        
        opps = []
        for i, row in enumerate(rows):
            # Shorten description to fit PPT
            desc_text = (row[1][:110] + '...') if len(row[1]) > 110 else row[1]
            title_text = f"#{i+1} {row[0]}  —  Score: {row[3]:.1f}/100"
            desc_full = f"{desc_text} · {row[2]} conversation{'s' if row[2] != 1 else ''}"
            opps.append((title_text, desc_full, colors[i]))
            
    except Exception as e:
        print("Could not load from DB:", e)
        opps = []
    
    top_offset = 2.8
    for title, desc, color in opps:
        # Title
        tx = slide.shapes.add_textbox(Inches(6.8), Inches(top_offset), Inches(6.0), Inches(0.3))
        tf = tx.text_frame
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Description
        tx2 = slide.shapes.add_textbox(Inches(6.8), Inches(top_offset + 0.22), Inches(6.0), Inches(0.3))
        tf2 = tx2.text_frame
        tf2.margin_top = Inches(0)
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.name = 'Calibri'
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(180, 180, 180)
        
        top_offset += 0.55

    prs.save('Myntra_Mindsight_Presentation_New.pptx')
    print("Presentation created successfully.")

if __name__ == '__main__':
    main()
